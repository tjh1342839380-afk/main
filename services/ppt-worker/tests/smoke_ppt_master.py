from __future__ import annotations

import json
import os
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation


def run(args: list[str]) -> None:
    result = subprocess.run(args, capture_output=True, text=True, encoding='utf-8', errors='replace')
    if result.returncode != 0:
        raise RuntimeError(result.stderr or result.stdout)


def main() -> None:
    root = Path(os.environ['PPT_MASTER_ROOT']).resolve()
    scripts = root / 'skills' / 'ppt-master' / 'scripts'
    with tempfile.TemporaryDirectory(prefix='ppt-master-smoke-') as temp:
        work = Path(temp)
        source = work / 'template.pptx'
        library = work / 'template.slide_library.json'
        plan_path = work / 'fill_plan.json'
        report = work / 'check_report.json'
        output = work / 'result_20000101_000000.pptx'

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = 'Original title'
        slide.placeholders[1].text = 'Original body'
        presentation.save(source)

        run([sys.executable, str(scripts / 'template_fill_pptx.py'), 'analyze', str(source), '-o', str(library)])
        run([sys.executable, str(scripts / 'template_fill_pptx.py'), 'scaffold', str(library), '-o', str(plan_path), '--slides', '1'])
        plan = json.loads(plan_path.read_text(encoding='utf-8'))
        plan['status'] = 'confirmed'
        plan['slides'][0]['replacements'][0]['text'] = 'Filled by PPT Master'
        plan_path.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding='utf-8')

        run([sys.executable, str(scripts / 'template_fill_pptx.py'), 'check-plan', str(library), str(plan_path), '-o', str(report)])
        run([sys.executable, str(scripts / 'template_fill_pptx.py'), 'apply', str(source), str(plan_path), '-o', str(output)])

        with zipfile.ZipFile(output) as archive:
            assert 'ppt/presentation.xml' in archive.namelist()
        assert json.loads(report.read_text(encoding='utf-8'))['summary']['error'] == 0
        print('PPT Master analyze -> scaffold -> check-plan -> apply: OK')


if __name__ == '__main__':
    main()
